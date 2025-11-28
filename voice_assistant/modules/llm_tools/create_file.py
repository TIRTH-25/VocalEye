# create_file.py — Improved: aggressive sanitization + robust PPTX pagination (no overflow)

import google.generativeai as genai
from config.settings import GEMINI_API_KEY, CURRENT_OS, DEFAULT_SAVE_PATH
import os, re, subprocess, hashlib

from reportlab.lib.pagesizes import A4
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, ListFlowable, ListItem, PageBreak
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.lib.enums import TA_LEFT
from reportlab.lib.colors import HexColor

from docx import Document
from docx.shared import Pt, RGBColor
from docx.enum.style import WD_STYLE_TYPE

from pptx import Presentation
from pptx.util import Inches, Pt as PPTPt
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.dml.color import RGBColor as PPTXRGB


# ---------------------------
# Palettes (unchanged)
# ---------------------------
PALETTES = {
    "deep_sapphire": {
        "bg": (6, 14, 38),
        "accent": (82, 145, 255),
        "accent2": (255, 188, 66),
        "text_main": (235, 244, 255),
        "text_soft": (178, 196, 220),
    },
    "slate_modern": {
        "bg": (20, 22, 25),
        "accent": (255, 131, 46),
        "accent2": (0, 165, 197),
        "text_main": (240, 244, 248),
        "text_soft": (168, 176, 188),
    },
    "minimal_white": {
        "bg": (250, 250, 250),
        "accent": (0, 135, 102),
        "accent2": (0, 90, 140),
        "text_main": (20, 28, 34),
        "text_soft": (108, 118, 128),
    }
}


def pick_palette(key: str):
    h = hashlib.sha256(key.encode()).hexdigest()
    keys = list(PALETTES.keys())
    return PALETTES[keys[int(h[:8], 16) % len(keys)]]


# ---------------------------
# Gemini + I/O helpers
# ---------------------------
def _ensure_genai():
    if not GEMINI_API_KEY:
        raise RuntimeError("Gemini API key missing.")
    genai.configure(api_key=GEMINI_API_KEY)


def open_file(filepath):
    try:
        if CURRENT_OS.startswith("win"):
            os.startfile(filepath)
        elif CURRENT_OS.startswith("darwin"):
            subprocess.run(["open", filepath])
        else:
            subprocess.run(["xdg-open", filepath])
    except:
        pass


# ---------------------------
# Aggressive sanitization + parser
# ---------------------------
def _sanitize_inline_markup(s: str) -> str:
    """
    Remove any leftover markup like **bold**, *italic*, `code`, __under__, and stray asterisks.
    Replace common markdown emphasis with plain text (keep content, remove markers).
    """
    if not s:
        return s
    # Remove code ticks
    s = re.sub(r"`+(.+?)`+", r"\1", s)
    # Bold and italic variants
    s = re.sub(r"\*\*(.+?)\*\*", r"\1", s)
    s = re.sub(r"__(.+?)__", r"\1", s)
    s = re.sub(r"\*(.+?)\*", r"\1", s)
    s = re.sub(r"_(.+?)_", r"\1", s)
    # Remove any remaining stray asterisks/underscores/backticks
    s = re.sub(r"[*_`]+", "", s)
    # Trim
    return s.strip()


def parse_txt_structure(text: str):
    """
    Parse markdown-like source into blocks while aggressively stripping markup characters.
    Returns list of {"type": ..., "text": ...}
    """
    out = []
    if text is None:
        return out
    for raw_line in text.splitlines():
        line = raw_line.rstrip()
        if not line:
            continue
        stripped = line.lstrip()
        if stripped.startswith("# "):
            out.append({"type": "h1", "text": _sanitize_inline_markup(stripped[2:])})
        elif stripped.startswith("## "):
            out.append({"type": "h2", "text": _sanitize_inline_markup(stripped[3:])})
        elif stripped.startswith("### "):
            out.append({"type": "h3", "text": _sanitize_inline_markup(stripped[4:])})
        elif stripped.startswith("- "):
            out.append({"type": "bullet", "text": _sanitize_inline_markup(stripped[2:])})
        else:
            out.append({"type": "p", "text": _sanitize_inline_markup(stripped)})
    return out


# ---------------------------
# Entrypoint
# ---------------------------
def file_generator(topic, filetype="pdf", save_path=None):
    _ensure_genai()
    model = genai.GenerativeModel("gemini-2.0-flash")
    prompt = f"""
Generate detailed, well-structured professional content for the topic.

Use the following only for structure (these symbol markers must NOT appear in final files):
# Heading 1
## Heading 2
### Heading 3
- Bullet point
Plain paragraph text

Topic: {topic}
Please produce several sections with headings, subheadings, and bullets so files are rich and detailed.
"""
    resp = model.generate_content(prompt)
    raw = resp.text.strip()

    blocks = parse_txt_structure(raw)

    if not save_path:
        os.makedirs(DEFAULT_SAVE_PATH, exist_ok=True)
        save_path = os.path.join(DEFAULT_SAVE_PATH, topic.replace(" ", "_") + ".txt")

    with open(save_path, "w", encoding="utf-8") as f:
        f.write(raw)

    if filetype == "pdf":
        out = save_path.replace(".txt", ".pdf")
        txt_to_pdf(blocks, out, topic)
        return out
    elif filetype == "docx":
        out = save_path.replace(".txt", ".docx")
        txt_to_docx(blocks, out, topic)
        return out
    elif filetype == "pptx":
        out = save_path.replace(".txt", ".pptx")
        txt_to_pptx(blocks, out, topic)
        return out

    return save_path


# ---------------------------
# PDF helper (hex fix already applied)
# ---------------------------
def _rgb_hex(rgb_tuple):
    r, g, b = rgb_tuple
    return f"#{r:02x}{g:02x}{b:02x}"


def txt_to_pdf(blocks, pdf_file, title):
    palette = pick_palette(title)
    accent = palette["accent"]
    accent2 = palette["accent2"]
    text_main = palette["text_main"]
    text_soft = palette["text_soft"]

    doc = SimpleDocTemplate(pdf_file, pagesize=A4, leftMargin=36, rightMargin=36, topMargin=48, bottomMargin=36)
    styles = getSampleStyleSheet()
    styles.add(ParagraphStyle(name="TitleX", fontSize=28, leading=32, textColor=HexColor(_rgb_hex(text_main)), spaceAfter=10))
    styles.add(ParagraphStyle(name="H1X", fontSize=20, leading=26, textColor=HexColor(_rgb_hex(accent)), spaceBefore=16, spaceAfter=6))
    styles.add(ParagraphStyle(name="H2X", fontSize=16, leading=22, textColor=HexColor(_rgb_hex(accent2)), spaceBefore=10, spaceAfter=4))
    styles.add(ParagraphStyle(name="H3X", fontSize=13, leading=18, textColor=HexColor(_rgb_hex(text_main)), spaceBefore=6, spaceAfter=2))
    styles.add(ParagraphStyle(name="BodyX", fontSize=11, leading=16, textColor=HexColor(_rgb_hex(text_soft))))

    story = []
    story.append(Paragraph(title, styles["TitleX"]))
    story.append(Spacer(1, 8))
    story.append(Paragraph(f'<para backColor="{_rgb_hex(accent)}"> </para>', ParagraphStyle(name='Bar', fontSize=1, leading=4)))
    story.append(Spacer(1, 10))

    bullets = []

    def flush_bullets():
        nonlocal bullets
        if bullets:
            story.append(ListFlowable([ListItem(Paragraph(b, styles["BodyX"])) for b in bullets], bulletType="bullet"))
            bullets = []
            story.append(Spacer(1, 6))

    for b in blocks:
        t, tx = b["type"], b["text"]
        if t == "h1":
            flush_bullets()
            story.append(PageBreak())
            story.append(Paragraph(tx, styles["H1X"]))
        elif t == "h2":
            flush_bullets()
            story.append(Paragraph(tx, styles["H2X"]))
        elif t == "h3":
            flush_bullets()
            story.append(Paragraph(tx, styles["H3X"]))
        elif t == "bullet":
            bullets.append(tx)
        else:
            flush_bullets()
            story.append(Paragraph(tx, styles["BodyX"]))
            story.append(Spacer(1, 4))

    flush_bullets()
    doc.build(story)
    open_file(pdf_file)


# ---------------------------
# DOCX generator (sanitized content so no ** remain)
# ---------------------------
def txt_to_docx(blocks, docx_file, title):
    palette = pick_palette(title)
    accent = palette["accent"]
    accent2 = palette["accent2"]
    text_main = palette["text_main"]
    text_soft = palette["text_soft"]

    doc = Document()
    styles = doc.styles

    # Title
    s = styles.add_style("TitleX", WD_STYLE_TYPE.PARAGRAPH)
    s.font.name = "Calibri"; s.font.size = Pt(32); s.font.bold = True; s.font.color.rgb = RGBColor(*text_main)
    s = styles.add_style("H1X", WD_STYLE_TYPE.PARAGRAPH)
    s.font.name = "Calibri"; s.font.size = Pt(22); s.font.bold = True; s.font.color.rgb = RGBColor(*accent)
    s = styles.add_style("H2X", WD_STYLE_TYPE.PARAGRAPH)
    s.font.name = "Calibri"; s.font.size = Pt(17); s.font.bold = True; s.font.color.rgb = RGBColor(*accent2)
    s = styles.add_style("H3X", WD_STYLE_TYPE.PARAGRAPH)
    s.font.name = "Calibri"; s.font.size = Pt(14); s.font.bold = True; s.font.color.rgb = RGBColor(*text_main)
    s = styles.add_style("BodyX", WD_STYLE_TYPE.PARAGRAPH)
    s.font.name = "Calibri"; s.font.size = Pt(12); s.font.color.rgb = RGBColor(*text_soft)

    doc.add_paragraph(title, style="TitleX")
    doc.add_paragraph("")

    for b in blocks:
        t, tx = b["type"], b["text"]
        if t == "h1":
            doc.add_page_break()
            doc.add_paragraph(tx, style="H1X")
        elif t == "h2":
            doc.add_paragraph(tx, style="H2X")
        elif t == "h3":
            doc.add_paragraph(tx, style="H3X")
        elif t == "bullet":
            p = doc.add_paragraph(tx, style="List Bullet")
            for r in p.runs:
                r.font.name = "Calibri"; r.font.size = Pt(12); r.font.color.rgb = RGBColor(*text_soft)
        else:
            p = doc.add_paragraph(tx, style="BodyX")
            p.space_after = Pt(4)

    doc.save(docx_file)
    open_file(docx_file)


# ---------------------------
# PPTX generator — robust pagination at write-time
# ---------------------------
def _estimate_chars_per_line(width_in):
    # More conservative: assume ~10 chars per inch (smaller than before)
    return max(30, int(width_in * 10))


def _lines_for_paragraph(text, chars_per_line):
    if not text:
        return [""]
    words = text.split()
    lines = []
    cur = words[0]
    for w in words[1:]:
        if len(cur) + 1 + len(w) <= chars_per_line:
            cur += " " + w
        else:
            lines.append(cur)
            cur = w
    lines.append(cur)
    return lines


def _slide_new(prs, bg, accent, left_strip=True):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    # background
    fill = slide.background.fill
    fill.solid()
    r, g, b = bg
    fill.fore_color.rgb = PPTXRGB(r, g, b)
    # left strip
    if left_strip:
        shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0), Inches(0), Inches(0.36), prs.slide_height)
        shfill = shape.fill
        shfill.solid()
        shfill.fore_color.rgb = PPTXRGB(*accent)
        shape.line.fill.background()
    return slide

from pptx.enum.text import MSO_AUTO_SIZE
import re
from pptx.dml.color import RGBColor as PPTXRGB
from pptx.util import Inches, Pt as PPTPt
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx import Presentation

def txt_to_pptx(blocks, pptx_file, title):
    palette = pick_palette(title)
    bg = palette["bg"]
    accent = palette["accent"]
    accent2 = palette["accent2"]
    text_main = palette["text_main"]
    text_soft = palette["text_soft"]

    prs = Presentation()

    # ------------------------
    # Cover slide
    # ------------------------
    cover = prs.slides.add_slide(prs.slide_layouts[6])
    f = cover.background.fill; f.solid(); f.fore_color.rgb = PPTXRGB(*bg)

    # top accent bar
    rect = cover.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0), Inches(0),
                                  prs.slide_width, Inches(0.5))
    rect.fill.solid(); rect.fill.fore_color.rgb = PPTXRGB(*accent)
    rect.line.fill.background()

    # cover title
    tb = cover.shapes.add_textbox(Inches(0.8), Inches(1.3), Inches(9), Inches(2))
    tf = tb.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.NONE
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = PPTPt(40)
    p.font.color.rgb = PPTXRGB(*text_main)
    p.font.bold = True

    # -------------------------------------
    # Conservative content layout parameters
    # -------------------------------------
    TEXT_BOX_LEFT = Inches(0.8)
    TEXT_BOX_WIDTH = Inches(7.5)   # narrower, safe
    TEXT_BOX_TOP = Inches(1.0)
    TEXT_BOX_HEIGHT = Inches(5.0)

    # line length safety
    CHAR_PER_LINE = 55

    # line height
    FONT_PT = 16
    LINE_HEIGHT_PT = FONT_PT * 1.3
    MAX_LINES = int((TEXT_BOX_HEIGHT.inches * 72) // LINE_HEIGHT_PT)

    # -----------------------------
    # Utility: new content slide
    # -----------------------------
    def new_slide():
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        f = slide.background.fill; f.solid(); f.fore_color.rgb = PPTXRGB(*bg)

        # left accent bar
        rect = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE,
                                      Inches(0), Inches(0), Inches(0.35), prs.slide_height)
        rect.fill.solid(); rect.fill.fore_color.rgb = PPTXRGB(*accent)
        rect.line.fill.background()

        box = slide.shapes.add_textbox(TEXT_BOX_LEFT, TEXT_BOX_TOP,
                                       TEXT_BOX_WIDTH, TEXT_BOX_HEIGHT)
        tf = box.text_frame
        tf.clear()
        tf.word_wrap = True
        tf.auto_size = MSO_AUTO_SIZE.NONE
        return slide, tf

    # active slide
    slide, tf = new_slide()
    used_lines = 0

    # --------------------------
    # Helper: ensure new slide
    # --------------------------
    def ensure_space(lines_needed):
        nonlocal slide, tf, used_lines
        if lines_needed + used_lines > MAX_LINES:
            slide, tf = new_slide()
            used_lines = 0
        return tf

    # -------------------------------------
    # Helper: soft-wrap long unbroken words
    # -------------------------------------
    # Use a normal string for replacement so '\u00AD' is a valid escape.
    SOFT_HYPHEN_REPL = "\\1\u00AD"
    def soften(text):
        # insert soft hyphen after every 15 consecutive word characters to allow breaking
        return re.sub(r"(\w{15})", SOFT_HYPHEN_REPL, text)

    # ---------------------------
    # Main content loop
    # ---------------------------
    for blk in blocks:
        t, tx = blk["type"], soften(blk["text"])

        # estimate line count (conservative)
        words = tx.split()
        lines = []
        cur = ""
        for w in words:
            if len((cur + " " + w).strip()) <= CHAR_PER_LINE:
                cur = (cur + " " + w).strip()
            else:
                lines.append(cur)
                cur = w
        if cur:
            lines.append(cur)
        needed = max(1, len(lines))

        # Heading slides
        if t == "h1":
            slide, tf = new_slide()
            tb = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(2))
            ttf = tb.text_frame; ttf.clear()
            p = ttf.paragraphs[0]
            p.text = tx
            p.font.size = PPTPt(32)
            p.font.color.rgb = PPTXRGB(*accent)
            p.font.bold = True

            slide, tf = new_slide()
            used_lines = 0
            continue

        if t == "h2":
            slide, tf = new_slide()
            tb = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(1))
            ttf = tb.text_frame; ttf.clear()
            p = ttf.paragraphs[0]
            p.text = tx
            p.font.size = PPTPt(26)
            p.font.color.rgb = PPTXRGB(*text_main)
            p.font.bold = True

            box = slide.shapes.add_textbox(TEXT_BOX_LEFT, Inches(1.8),
                                           TEXT_BOX_WIDTH, TEXT_BOX_HEIGHT)
            tf = box.text_frame
            tf.clear()
            tf.word_wrap = True
            tf.auto_size = MSO_AUTO_SIZE.NONE

            used_lines = 0
            continue

        # normal content
        # ensure space before writing
        tf = ensure_space(needed)

        for ln in lines:
            p = tf.add_paragraph()
            p.text = ln
            p.font.size = PPTPt(FONT_PT)
            p.font.name = "Calibri"
            p.font.color.rgb = PPTXRGB(*text_soft)
            p.level = 0
        used_lines += needed

    prs.save(pptx_file)
    open_file(pptx_file)



